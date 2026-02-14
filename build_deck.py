"""
Build 'AI in Biopharma' Bain POV PowerPoint deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Bain palette ──
BAIN_RED = RGBColor(0xCC, 0x00, 0x00)
BAIN_DARK = RGBColor(0x8B, 0x00, 0x00)
BAIN_LIGHT = RGBColor(0xE6, 0x45, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
TABLE_HEADER_BG = RGBColor(0xCC, 0x00, 0x00)
TABLE_ALT_BG = RGBColor(0xFA, 0xFA, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ──

def add_red_bar(slide, top=0, height=Inches(0.06)):
    """Add a thin Bain red bar across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BAIN_RED
    shape.line.fill.background()

def add_footer(slide, text="Bain & Company  |  Confidential  |  February 2026"):
    left = Inches(0.5)
    top = prs.slide_height - Inches(0.45)
    txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = MED_GRAY
    p.font.name = "Calibri"

def add_page_number(slide, num):
    left = prs.slide_width - Inches(1)
    top = prs.slide_height - Inches(0.45)
    txBox = slide.shapes.add_textbox(left, top, Inches(0.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(8)
    p.font.color.rgb = MED_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT

def set_cell_text(cell, text, font_size=10, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_table(slide, rows_data, col_widths, left, top, row_height=Inches(0.45)):
    n_rows = len(rows_data)
    n_cols = len(col_widths)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, row_height * n_rows)
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            is_header = (ri == 0)
            set_cell_text(
                cell, val,
                font_size=9 if not is_header else 9,
                bold=is_header,
                color=WHITE if is_header else BLACK,
            )
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
    return table

def title_slide(slide, title_text, subtitle_text):
    """Full red background title slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BAIN_RED
    bg.line.fill.background()
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.3), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.LEFT
    # Bottom branding
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(6), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Bain & Company  |  February 2026  |  Confidential"
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(0xFF, 0xBD, 0xBD)
    p3.font.name = "Calibri"

def content_slide(slide, title_text, subtitle_text=None):
    """Standard content slide with red bar, title, optional subtitle."""
    add_red_bar(slide)
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(11.9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.font.name = "Calibri"
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(11.9), Inches(0.45))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.italic = True
        p2.font.color.rgb = BAIN_RED
        p2.font.name = "Calibri"

def add_body_text(slide, text, left=0.7, top=1.5, width=11.9, height=5.2, font_size=12):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        if line.startswith("   "):
            p.level = 1
            p.font.size = Pt(font_size - 1)
            p.font.color.rgb = GRAY

def add_bullets(slide, bullets, left=0.7, top=1.5, width=11.9, height=5.2, font_size=12):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level, bold, color, sz) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz if sz else font_size)
        p.font.bold = bold
        p.font.color.rgb = color if color else BLACK
        p.font.name = "Calibri"
        p.level = level
        p.space_after = Pt(3)
    return txBox

# ════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
title_slide(
    slide,
    "AI in Biopharma: The Trillion-Dollar Blind Spot",
    "Why the most data-rich industry on earth is the least AI-native \u2014\nand three unconventional plays to change the game"
)

# ════════════════════════════════════════════════════════
# SLIDE 2 — THE PROVOCATION
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "The Provocation", '"Pharma has more data than any industry in the world. And less intelligence from it."')

bullets = [
    ("Pharma generates ~2 exabytes of clinical, genomic, and commercial data annually \u2014 yet fewer than 20% of enterprises have scaled AI in any meaningful way", 0, False, BLACK, 12),
    ("R&D productivity has declined ~50% over the last decade; Eroom's Law persists despite $250B+ in annual R&D spend globally", 0, False, BLACK, 12),
    ("The industry is trapped in what Bain calls the 'micro-productivity trap' \u2014 hundreds of disconnected pilots, impressive demos, zero enterprise value", 0, False, BLACK, 12),
    ("85% of providers believe AI will transform diagnosis/treatment in 3\u20135 years; 25% of ChatGPT users already submit healthcare questions weekly", 0, False, BLACK, 12),
    ("Utah has authorized AI to legally prescribe routine refills. OpenAI launched ChatGPT Health to analyze medical records.", 0, False, BLACK, 12),
    ("", 0, False, BLACK, 12),
    ("The dire truth: While pharma debates governance frameworks and runs chatbot pilots, tech-native insurgents are building the capabilities that will disintermediate incumbents across the entire value chain.", 0, True, BAIN_RED, 13),
]
add_bullets(slide, bullets, top=1.5)
add_footer(slide)
add_page_number(slide, 2)

# ════════════════════════════════════════════════════════
# SLIDE 3 — THE DIAGNOSIS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "The Diagnosis: Three Root Causes of Pharma's AI Paralysis")

rows = [
    ["Root Cause", "What We See", "Why It Matters"],
    [
        "1. AI treated as IT project,\nnot business transformation",
        "CIOs adopt 'wait and see'; CDOs funded without P&L ownership; grassroots experimentation fragments into silos",
        "Winners achieving 10\u201325% EBIT improvement through zero-based process redesign; laggards automate individual tasks and return to status quo"
    ],
    [
        "2. Data without meaning:\nno semantic layer",
        "Massive data lakes, but no governed business definitions; every function defines 'revenue,' 'patient,' 'signal' differently; AI agents hallucinate meaning",
        "AI produces confident but wrong answers (context poisoning); trust erodes before value materializes; different agents give different answers to the same question"
    ],
    [
        "3. Value chain fragmentation:\nkingdoms, not an enterprise",
        "R&D, Manufacturing, Commercial, and Medical Affairs operate disconnected AI agendas; ~80% of ERP transformations miss goals",
        "No end-to-end AI architecture; cross-functional use cases (evidence \u2192 launch \u2192 commercial) remain unrealized; duplication of work across assets and TAs"
    ],
]
add_table(slide, rows, [Inches(3.2), Inches(4.3), Inches(4.4)], Inches(0.7), Inches(1.3), row_height=Inches(1.35))
add_footer(slide)
add_page_number(slide, 3)

# ════════════════════════════════════════════════════════
# SLIDE 4 — THE MISALLOCATION
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Where Pharma AI Spending Goes Today", '"A thousand points of light, none bright enough to see by"')

bullets = [
    ("70%+ of pharma AI investment goes to: copilots for search/summarization, literature reviews, chatbot prototypes", 0, False, BLACK, 12),
    ("<10% goes to process redesign, operating model change, or value-chain-spanning systems", 0, False, BLACK, 12),
    ("Result: Localized 5\u201310% efficiency gains that never scale beyond the pilot team", 0, False, BLACK, 12),
    ("", 0, False, BLACK, 10),
    ("The comparison:", 0, True, BAIN_RED, 13),
    ("Banking: One bank turned customer insight \u2192 campaign from 60\u2013100 days to 1 day; 40 employees and 10 handoffs reduced to 4\u20135 with zero handoffs", 1, False, GRAY, 11),
    ("Automotive: OEMs that act on AI-first sales face up to 20% top-line increase; laggards risk losing 15%", 1, False, GRAY, 11),
    ("Pharma: Average 1.5:1 ROI on pilot use cases; no scaled transformations in the public domain", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 10),
    ("The wake-up call:", 0, True, BAIN_RED, 13),
    ("AI at the point of care is already here. If pharma evidence is not discoverable, interpretable, and citable by AI tools, it will be deprioritized or absent in AI-driven clinical decisions. Not engaging risks your evidence being misinterpreted, deprioritized, or absent.", 0, False, BLACK, 12),
]
add_bullets(slide, bullets, top=1.5)
add_footer(slide)
add_page_number(slide, 4)

# ════════════════════════════════════════════════════════
# SLIDE 5 — THE STRUCTURAL SHIFT
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "The Industry Is Missing the Structural Shift", '"Point of care decision-making is being reshaped by AI \u2014 and pharma isn\u2019t in the room"')

bullets = [
    ("AI is not a tool pharma uses. AI is becoming the medium through which treatment decisions are made.", 0, True, BLACK, 13),
    ("", 0, False, BLACK, 8),
    ("The CRO market (~$80B) is at an inflection point:", 0, True, BAIN_RED, 12),
    ("AI early adopters (IQVIA+NVIDIA, PPD+OpenAI, Parexel+Palantir) are fundamentally changing trial delivery", 1, False, GRAY, 11),
    ("CROs who invest: faster site selection, 90% reduction in SOP search time, automated safety signal detection, CRA co-pilots", 1, False, GRAY, 11),
    ("CROs who don't: margin compression, talent attrition, customer defection", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 8),
    ("Yet the industry is moving backward:", 0, True, BAIN_RED, 12),
    ("54% of large pharma is reducing CRO outsourcing; shifting to FSP models for cost savings \u2014 not investing in AI-native trial delivery", 1, False, GRAY, 11),
    ("R&D spend growth moderating to low-single digits; trial starts declining, especially post-Phase I", 1, False, GRAY, 11),
    ("Biotech funding recovering but concentrated in 'top programs' \u2014 the long tail is starving", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 8),
    ("The window to act is 12\u201318 months. After that, the AI-mediated point of care will be established \u2014 with or without your evidence in it.", 0, True, BAIN_DARK, 12),
]
add_bullets(slide, bullets, top=1.5)
add_footer(slide)
add_page_number(slide, 5)

# ════════════════════════════════════════════════════════
# SLIDE 6 — TRANSITION: THREE PLAYS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Three Plays Nobody in the Industry Is Making"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p2 = tf.add_paragraph()
p2.text = '"Stop thinking about AI use cases. Start thinking about AI-native business models."'
p2.font.size = Pt(14)
p2.font.italic = True
p2.font.color.rgb = BAIN_LIGHT
p2.font.name = "Calibri"

# Three boxes
plays = [
    ("PLAY 1", "The Evidence\nIntelligence Engine", "Own the AI-mediated\npoint of care before\nsomeone else does"),
    ("PLAY 2", "The Autonomous\nClinical Machine", "Build the self-driving\nclinical trial \u2014 not\nAI-assisted, AI-native"),
    ("PLAY 3", "The Pharma\nSemantic Enterprise", "Create the industry's\nfirst AI-native\noperating system"),
]
for i, (label, title, desc) in enumerate(plays):
    left = Inches(1 + i * 3.9)
    top = Inches(2.8)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.5), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = BAIN_RED
    box.line.fill.background()
    # Label
    t = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), Inches(2.9), Inches(0.35))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xBD, 0xBD)
    p.font.name = "Calibri"
    # Title
    t2 = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.7), Inches(2.9), Inches(1.2))
    tf2 = t2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"
    # Desc
    t3 = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(2.0), Inches(2.9), Inches(1.2))
    tf3 = t3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0xFF, 0xDD, 0xDD)
    p3.font.name = "Calibri"

# ════════════════════════════════════════════════════════
# SLIDE 7 — PLAY 1 OVERVIEW
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Play 1: The Evidence Intelligence Engine", "Own the AI-mediated point of care before someone else does")

rows = [
    ["Function", "Opportunity", "Capability Required", "Value Potential"],
    [
        "R&D",
        "Design endpoints, trial outcomes, and evidence packages specifically for AI interpretation \u2014 multimodal biomarkers, longitudinal RWE, machine-readable study reports",
        "AI-native evidence design standards; partnership with CDS tool developers; RWE designed for AI interpretation with advanced endpoints",
        "15\u201320% reduction in post-launch evidence gaps; faster HTA approval; fewer protocol amendments"
    ],
    [
        "Commercial",
        "Build CDS API layers ensuring your evidence is discoverable and cited by AI tools clinicians use at POC; monitor what clinicians ask AI and adapt messaging in real-time",
        "Generative Experience Optimization (GXO) for clinical content; real-time Q&A monitoring infrastructure; CDS-ready content library",
        "Up to 20% top-line impact for movers vs. 15% loss for laggards; real-time competitive intelligence from AI POC queries"
    ],
    [
        "Medical Affairs",
        "Deploy agentic 'Rubik's Cube' advisory boards \u2014 AI personas (patient, payer, physician, regulator, caregiver) stress-test evidence packages before real engagement",
        "Multi-agent simulation platform; stakeholder persona models trained on real survey data and regulatory precedent",
        "3x faster evidence strategy iteration; 60\u201370% reduction in advisory board costs; unprecedented multi-stakeholder balancing"
    ],
]
add_table(slide, rows, [Inches(1.3), Inches(3.4), Inches(3.4), Inches(3.4)], Inches(0.5), Inches(1.4), row_height=Inches(1.3))
add_footer(slide)
add_page_number(slide, 7)

# ════════════════════════════════════════════════════════
# SLIDE 8 — PLAY 1 DEEP DIVE: AGENTIC ADVISORY BOARD
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Play 1 Deep Dive: The Agentic Advisory Board", "From episodic advisory boards to continuous AI-simulated stakeholder feedback")

bullets = [
    ("Today: Pharma convenes advisory boards 2\u20133x per asset lifecycle", 0, True, BLACK, 13),
    ("Slow (months to organize), expensive ($50K\u2013$200K each), limited perspectives, episodic insights", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Tomorrow: AI agents representing key stakeholder personas interrogate evidence packages on demand", 0, True, BLACK, 13),
    ("Physician Agent: Evaluates clinical relevance, endpoint significance, treatment positioning", 1, False, GRAY, 11),
    ("Payer Agent: Assesses reimbursement likelihood, cost-effectiveness, population evidence (95.6% accuracy demonstrated \u2014 RISA Labs)", 1, False, GRAY, 11),
    ("Patient Agent: Flags patient-relevant gaps, adherence barriers, PRO alignment", 1, False, GRAY, 11),
    ("Regulator Agent: Predicts health authority queries, identifies evidence gaps for label negotiations", 1, False, GRAY, 11),
    ("Caregiver Agent: Surfaces real-world care burden and unmet support needs", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Proof points from analogous deployments:", 0, True, BAIN_RED, 12),
    ("Tsinghua multi-agent hospital simulation: ~93% accuracy on clinical reasoning", 1, False, GRAY, 11),
    ("Oxford/Hull regulator-manufacturer simulation: qualitative insight into compliance vs. innovation tradeoffs", 1, False, GRAY, 11),
    ("Commercial solutions emerging: PosterosAI (scientific intelligence), others (provider simulation, rep training)", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Where to start: Prototype on one live asset; 3\u20135 evidence questions on existing data; validate against real advisory board output", 0, True, BAIN_DARK, 12),
]
add_bullets(slide, bullets, top=1.4)
add_footer(slide)
add_page_number(slide, 8)

# ════════════════════════════════════════════════════════
# SLIDE 9 — PLAY 2 OVERVIEW
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Play 2: The Autonomous Clinical Machine", "Build the self-driving clinical trial \u2014 not AI-assisted, AI-native")

rows = [
    ["Function", "Opportunity", "Capability Required", "Value Potential"],
    [
        "R&D / Clinical",
        "AI-native protocol design (predictive amendment avoidance); ML site selection; continuous risk-based monitoring; LLM signal extraction from social/EHR for early safety detection",
        "Integrated AI clinical platform; real-time data layer; CRA co-pilot; NLP pipeline for PV signal detection from unstructured data (Reddit, EHR notes)",
        "30\u201340% reduction in trial timelines; 25% lower cost per patient; earlier safety signal detection"
    ],
    [
        "Manufacturing",
        "Predictive clinical supply planning across sites/geographies; AI-optimized scale-up from clinical to commercial; demand signal integration from enrollment data",
        "Digital twin of clinical supply chain; ML demand forecasting; semantic integration with R&D data",
        "20\u201330% reduction in clinical supply waste; faster commercial manufacturing readiness"
    ],
    [
        "Commercial",
        "Real-time enrollment intelligence feeds pre-launch planning; AI-generated market access dossiers from trial data",
        "Cross-functional data pipeline from R&D to Commercial; automated evidence synthesis tools",
        "6\u201312 months earlier commercial preparation; stronger launch evidence packages"
    ],
    [
        "Medical Affairs",
        "Automated safety intake, case reporting, aggregate reporting; AI-powered pharmacovigilance from unstructured data sources",
        "LLM safety signal extraction; knowledge graph for drug side effects; automated FAERS integration",
        "200K+ hours/year saved in safety ops; earlier PV signals; Spearman ~0.42 correlation between social media signals and FAERS data"
    ],
]
add_table(slide, rows, [Inches(1.3), Inches(3.2), Inches(3.5), Inches(3.5)], Inches(0.5), Inches(1.4), row_height=Inches(1.15))
add_footer(slide)
add_page_number(slide, 9)

# ════════════════════════════════════════════════════════
# SLIDE 10 — PLAY 2 DEEP DIVE: CRO DISRUPTION
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Play 2 Deep Dive: The CRO Disruption Angle", "AI will create winners and losers in the $80B CRO market \u2014 pharma must be on the right side")

bullets_left = [
    ("CROs at the forefront of AI investment:", 0, True, BLACK, 13),
    ("IQVIA + NVIDIA: Healthcare-grade agentic AI for protocol adherence, data accuracy, oversight across trials", 1, False, GRAY, 11),
    ("PPD (Thermo Fisher) + OpenAI: AI-powered drug development platform from early-stage to commercialization", 1, False, GRAY, 11),
    ("Parexel + Palantir: Enhanced clinical data workflows; AI-native trial operations model", 1, False, GRAY, 11),
    ("Fortrea: AI in Clinical Development initiative; $150M gross savings target", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("AI across the full trial lifecycle:", 0, True, BLACK, 13),
    ("Trial Design: AI-supported protocol authoring, IB generation, literature review (esp. SLRs)", 1, False, GRAY, 11),
    ("Trial Start-up: AI/ML site selection and prioritization; site onboarding material generation", 1, False, GRAY, 11),
    ("Patient Recruitment: Patient identification, matching, engagement; modular content generation", 1, False, GRAY, 11),
    ("Data Collection: EDC database build automation; automated data checking/cleaning; CRA co-pilot", 1, False, GRAY, 11),
    ("Trial Close-out: Clinical Study Report generation; regulatory query response drafting", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("The unconventional play:", 0, True, BAIN_RED, 13),
    ("Don't just outsource trials \u2014 co-develop AI-native trial capabilities with CRO partners. Use AI to reduce scope of what needs to be outsourced while increasing quality of what is retained.", 0, False, BLACK, 12),
    ("Value creation levers: Labor optimization (workforce transformation), tech-enabled delivery, AI-driven site strategy, outcome-based pricing enabled by AI prediction confidence", 0, False, BLACK, 11),
]
add_bullets(slide, bullets_left, top=1.4)
add_footer(slide)
add_page_number(slide, 10)

# ════════════════════════════════════════════════════════
# SLIDE 11 — PLAY 3 OVERVIEW
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Play 3: The Pharma Semantic Enterprise", "Create the industry's first AI-native operating system across the value chain")

rows = [
    ["Function", "Opportunity", "Capability Required", "Value Potential"],
    [
        "R&D",
        "Unified definitions of 'patient,' 'endpoint,' 'signal,' 'evidence' across discovery, clinical, and regulatory \u2014 so AI agents can reason across the entire pipeline",
        "Enterprise knowledge graph; pharma-specific ontology (drug, trial, indication, patient, evidence); governed semantic registry",
        "Eliminate 40\u201360% of reconciliation work between functions; cross-asset evidence reuse; faster IEP creation"
    ],
    [
        "Manufacturing",
        "Semantic integration of manufacturing data with R&D and commercial \u2014 AI can trace from formulation to patient outcome to demand signal",
        "Manufacturing semantic model connected to R&D and commercial graphs; digital thread across CMC",
        "15\u201320% improvement in tech transfer speed; predictive quality management"
    ],
    [
        "Commercial",
        "One definition of 'customer,' 'HCP engagement,' 'territory,' 'sales' across field force, medical, and marketing \u2014 enabling AI-led segmentation and omnichannel orchestration",
        "Commercial semantic model; AI-led account segmentation engine; unified C360 platform",
        "25\u201330% improvement in field force productivity; hyper-personalized HCP engagement"
    ],
    [
        "Medical Affairs",
        "Governed, versioned definitions of evidence quality, signal strength, and regulatory classification \u2014 so AI agents in safety, med info, and publications work from same truth",
        "Medical affairs knowledge graph; governed evidence taxonomy; audit-ready lineage and provenance",
        "Eliminate 'shadow AI' risk in regulated functions; audit-ready AI outputs; compliance by design"
    ],
]
add_table(slide, rows, [Inches(1.3), Inches(3.2), Inches(3.5), Inches(3.5)], Inches(0.5), Inches(1.4), row_height=Inches(1.15))
add_footer(slide)
add_page_number(slide, 11)

# ════════════════════════════════════════════════════════
# SLIDE 12 — PLAY 3 DEEP DIVE: SEMANTIC LAYER
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Play 3 Deep Dive: From Data Lake to Decision Engine", "The Semantic Layer Architecture for Pharma")

bullets = [
    ("The problem:", 0, True, BAIN_RED, 13),
    ("Pharma has spent billions on data lakes and warehouses. AI agents query raw schemas and guess the meaning.", 0, False, BLACK, 12),
    ("Different agents define 'revenue,' 'patient adherence,' 'site performance' differently. Leaders can't trust outputs.", 0, False, BLACK, 12),
    ("Context Poisoning: Outdated definitions lead to systematic errors. Context Confusion: Multiple valid meanings compete. Context Clash: Sources contradict each other.", 0, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("The solution: A governed semantic layer that defines business meaning once, makes it reusable everywhere", 0, True, BAIN_RED, 13),
    ("", 0, False, BLACK, 6),
    ("Architecture: Raw Data \u2192 Bronze (available) \u2192 Silver (reliable) \u2192 Gold (reportable) \u2192 Semantic Layer (understandable) \u2192 AI Agents / BI / Applications", 0, True, BLACK, 12),
    ("", 0, False, BLACK, 6),
    ("Five key components:", 0, True, BLACK, 13),
    ("1. Business Definitions: Canonical terms, metric definitions, calculation assumptions, business rules", 1, False, GRAY, 11),
    ("2. Metadata & Context: Data mapping, lineage & provenance, freshness & time context, access & entitlements", 1, False, GRAY, 11),
    ("3. Entities & Hierarchies: Entity models (patient, drug, site, trial), taxonomies, roll-up rules, aggregation logic", 1, False, GRAY, 11),
    ("4. Ontology & Relationships: Relationship types, join logic, constraints, causal/influence links", 1, False, GRAY, 11),
    ("5. Knowledge Graph: Connected entity graph, cross-domain reasoning paths, evidence & signal links", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Platform options: Palantir Foundry (operational), Stardog (standards-based), dbt + Looker (metric layer), or hybrid", 0, False, BLACK, 12),
    ("Proven at scale: Tyson, Lowe's, Uber, Cardinal Health, Klarna, Bosch, MLB \u2014 all using semantic layers for AI", 0, False, BLACK, 12),
]
add_bullets(slide, bullets, top=1.4, font_size=11)
add_footer(slide)
add_page_number(slide, 12)

# ════════════════════════════════════════════════════════
# SLIDE 13 — INTEGRATED VALUE MAP
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "The Integrated Value Map: AI Potential Across the Pharma Enterprise")

rows = [
    ["Value Chain Area", "Near-Term (12\u201318 mo)", "Medium-Term (18\u201336 mo)", "Transformational (36+ mo)", "Est. EBIT Impact"],
    [
        "R&D / Clinical",
        "AI protocol authoring; ML site selection; CRA co-pilot; LLM literature review; automated SOP search",
        "Autonomous trial monitoring; agentic evidence advisory; predictive enrollment; AI-native regulatory submissions",
        "Self-driving clinical trials; real-time evidence generation; AI-mediated regulatory engagement",
        "15\u201325%\ncost reduction"
    ],
    [
        "Manufacturing",
        "Predictive quality; AI-optimized clinical supply planning; process analytics co-pilot",
        "Digital twin manufacturing; semantic integration with R&D; AI-driven tech transfer",
        "Autonomous manufacturing operations; closed-loop quality management",
        "10\u201315%\ncost reduction"
    ],
    [
        "Commercial",
        "AI-led HCP segmentation; content personalization; GXO for CDS tools; competitive intelligence",
        "Omnichannel AI orchestration; real-time POC evidence delivery; AI-powered field coaching",
        "AI-mediated prescribing influence; autonomous commercial operations; 1:1 HCP engagement",
        "15\u201320%\nrevenue uplift"
    ],
    [
        "Medical Affairs",
        "SOP co-pilot (200K hrs saved); automated PV intake; signal detection from social/EHR",
        "Agentic stakeholder simulation; AI evidence gap analysis; automated publications",
        "Continuous RWE generation; autonomous safety monitoring; real-time medical intelligence",
        "10\u201315% cost\nreduction + risk"
    ],
    [
        "Enterprise / G&A",
        "AI-powered procurement; finance co-pilots; HR automation; agentic ERP workflows",
        "Touchless platform operations; multi-agent orchestration across SAP/Workday/Salesforce",
        "Fully autonomous support functions; AI as primary interface for enterprise systems",
        "20\u201330%\nG&A reduction"
    ],
]
add_table(slide, rows, [Inches(1.5), Inches(2.6), Inches(2.8), Inches(2.8), Inches(1.6)], Inches(0.5), Inches(1.3), row_height=Inches(1.05))
add_footer(slide)
add_page_number(slide, 13)

# ════════════════════════════════════════════════════════
# SLIDE 14 — WINNERS VS LAGGARDS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "What Separates Winners from Laggards", "The Bain AI Transformation Framework Applied to Pharma")

bullets = [
    ("1.  Top-down leadership", 0, True, BAIN_RED, 15),
    ("AI transformation starts in the C-suite, not the innovation lab. Integrate AI into strategy with ambition matched by executive ownership. Embed AI objectives into performance reviews, bonus structures, and promotion criteria. Shopify mandated AI as 'baseline expectation' \u2014 pharma CEOs must do the same.", 0, False, BLACK, 11),
    ("", 0, False, BLACK, 6),
    ("2.  Fewer, bigger bets", 0, True, BAIN_RED, 15),
    ("4\u20135 critical domains, not 200 pilots. In pharma: clinical development, evidence generation, commercial engagement, pharmacovigilance, and enterprise data. Each domain is a system of work with 20\u201340+ interrelated use cases. Transformation happens at the domain level, not the use case level.", 0, False, BLACK, 11),
    ("", 0, False, BLACK, 6),
    ("3.  Zero-based process redesign", 0, True, BAIN_RED, 15),
    ("Map point of departure \u2192 reimagine point of arrival with AI at the core. Don't automate broken workflows \u2014 build entirely new processes with AI at the center. The process redesign, not the technology, creates most of the value.", 0, False, BLACK, 11),
    ("", 0, False, BLACK, 6),
    ("4.  Operating model for continuous transformation", 0, True, BAIN_RED, 15),
    ("Small transformation team + business-owned solution teams. Two speeds: run the business and change the business simultaneously. Six critical areas: E2E process, solution team mobilization, data governance, scaling, adoption, business-tech partnership.", 0, False, BLACK, 11),
    ("", 0, False, BLACK, 6),
    ("5.  Semantic foundation first", 0, True, BAIN_RED, 15),
    ("Build the meaning layer before scaling agents. Context management is the enterprise capability that makes everything else work. Without it, every new AI agent multiplies inconsistency, not intelligence.", 0, False, BLACK, 11),
]
add_bullets(slide, bullets, top=1.4)
add_footer(slide)
add_page_number(slide, 14)

# ════════════════════════════════════════════════════════
# SLIDE 15 — COST OF INACTION
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11.3), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The Cost of Inaction"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"
p2 = tf.add_paragraph()
p2.text = '"The question is no longer how pharma will use AI. It\u2019s whether pharma will still control its own destiny in a world where AI mediates every clinical decision."'
p2.font.size = Pt(14)
p2.font.italic = True
p2.font.color.rgb = BAIN_LIGHT
p2.font.name = "Calibri"

bullets_data = [
    "AI leaders are delivering double the EBIT margins of their competitors",
    "85% of providers expect AI to transform treatment decisions in 3\u20135 years",
    "Biotech insurgents with AI-native operations need 60% fewer people and 40% less capital per trial",
    "CROs that don't invest in AI face margin compression, talent attrition, and customer defection",
    "Every month of delay compounds: competitors scale, data advantages widen, talent migrates, and the AI-mediated point of care evolves without your evidence in it",
]
top_y = Inches(2.8)
for i, text in enumerate(bullets_data):
    t = slide.shapes.add_textbox(Inches(1.3), top_y + Inches(i * 0.65), Inches(10.5), Inches(0.6))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "\u25B6  " + text
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

# Closing line
t_close = slide.shapes.add_textbox(Inches(1), Inches(6.1), Inches(11.3), Inches(0.6))
tf_c = t_close.text_frame
tf_c.word_wrap = True
p_c = tf_c.paragraphs[0]
p_c.text = "The winners won't have the most pilots. They'll have the fewest \u2014 and the biggest."
p_c.font.size = Pt(18)
p_c.font.bold = True
p_c.font.color.rgb = BAIN_LIGHT
p_c.font.name = "Calibri"
p_c.alignment = PP_ALIGN.CENTER

# ════════════════════════════════════════════════════════
# SLIDE 16 — RECOMMENDED NEXT STEPS
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Recommended Next Steps")

steps = [
    ("1", "Conduct an enterprise AI maturity diagnostic", "5-week sprint: ~30 interviews across functions (R&D, Commercial, Medical, Manufacturing, IT); synthesize maturity across 20+ subdimensions of data strategy, culture, governance, architecture, and analytics; benchmark vs. pharma and cross-industry best-in-class peers"),
    ("2", "Identify 4\u20135 high-value domains", "Build top-down value hypotheses with specific P&L targets for each domain; prioritize based on competitive advantage, data readiness, and organizational willingness; define go/no-go criteria for each"),
    ("3", "Stand up a semantic layer pilot", "Start in one high-priority domain (e.g., commercial data model or clinical trial ontology); define canonical terms, metrics, hierarchies; connect to 2\u20133 AI agent use cases; prove consistency and trust improvement"),
    ("4", "Prototype one agentic use case in evidence generation", "E.g., Patient-Aligned Clinical Evidence (PACE) diagnostic or Agentic Advisory Board; 3\u20135 evidence questions on existing data for one live asset; test usability, decision influence, and KPI impact"),
    ("5", "Design the transformation operating model", "Business-owned solution teams + small central transformation team; define governance, scaling mechanisms, adoption metrics; embed into existing stage-gate and IEP review processes; plan for team enablement and upskilling"),
]
for i, (num, title, desc) in enumerate(steps):
    top = Inches(1.3 + i * 1.15)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.05), Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BAIN_RED
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].text = num
    circ.text_frame.paragraphs[0].font.size = Pt(14)
    circ.text_frame.paragraphs[0].font.bold = True
    circ.text_frame.paragraphs[0].font.color.rgb = WHITE
    circ.text_frame.paragraphs[0].font.name = "Calibri"
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Title
    t = slide.shapes.add_textbox(Inches(1.3), top, Inches(11), Inches(0.35))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.font.name = "Calibri"
    # Description
    t2 = slide.shapes.add_textbox(Inches(1.3), top + Inches(0.35), Inches(11), Inches(0.7))
    tf2 = t2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"

add_footer(slide)
add_page_number(slide, 16)

# ════════════════════════════════════════════════════════
# SLIDE 17 — WHY BAIN
# ════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_slide(slide, "Why Bain Is the Right Partner for This Journey")

# Left column - credentials
bullets_left = [
    ("AI & Analytics Scale", 0, True, BAIN_RED, 14),
    ("3,000+ AI, Insights & Solutions client engagements globally", 1, False, GRAY, 11),
    ("420+ GenAI/Agentic AI use cases and solutions delivered worldwide", 1, False, GRAY, 11),
    ("250+ production solutions built worldwide via Vector", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Healthcare & Life Sciences Depth", 0, True, BAIN_RED, 14),
    ("20+ HLS AI projects in 2025 alone (Amgen, Novo, Pfizer, GEHC, IHH, and more)", 1, False, GRAY, 11),
    ("100+ data & analytics projects in healthcare", 1, False, GRAY, 11),
    ("85% of healthcare PE deal value advised", 1, False, GRAY, 11),
    ("Deep expertise across pharma value chain: R&D, Medical Affairs, Commercial, Manufacturing", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Strategic Alliances", 0, True, BAIN_RED, 14),
    ("OpenAI: 2.5+ year alliance; Sam Altman cited Bain at DevDay 2025", 1, False, GRAY, 11),
    ("Palantir: Enterprise-grade AI solutions at scale; HLS priority partnership", 1, False, GRAY, 11),
    ("Andrew Ng / AI Aspire: Accelerating AI transformation", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Talent", 0, True, BAIN_RED, 14),
    ("1,500+ deep analytics experts and practitioners", 1, False, GRAY, 11),
    ("900+ software engineers via partner network", 1, False, GRAY, 11),
    ("60+ strategic design, UI/UX and user research experts", 1, False, GRAY, 11),
    ("20+ senior practitioners (ex-CDO, CAO, data science professors)", 1, False, GRAY, 11),
    ("", 0, False, BLACK, 6),
    ("Results Focus", 0, True, BAIN_RED, 14),
    ("We don't work on 'Bain projects.' We partner in accelerating your work and results.", 0, False, BLACK, 12),
    ("We define and pursue success together \u2014 with scientific rigor in Results Delivery.", 0, False, BLACK, 12),
]
add_bullets(slide, bullets_left, top=1.3, font_size=11)
add_footer(slide)
add_page_number(slide, 17)

# ════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════
output_path = r"C:\Users\75565\OneDrive - Bain\Documents\GitHub\satty001\AI_in_Biopharma_Bain_POV.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
