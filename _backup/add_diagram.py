"""
Add a Clinical Development Machine of the Future value chain diagram
to the existing AI in Biopharma deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──
BAIN_RED     = RGBColor(0xCC, 0x00, 0x00)
BAIN_DARK    = RGBColor(0x8B, 0x00, 0x00)
BAIN_LIGHT   = RGBColor(0xE6, 0x45, 0x45)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x1A, 0x1A, 0x1A)
GRAY         = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xEE, 0xEE, 0xEE)
DARK_BG      = RGBColor(0x2D, 0x2D, 0x2D)
DARK_CARD    = RGBColor(0x3A, 0x3A, 0x3A)
ACCENT_GREEN = RGBColor(0x00, 0x96, 0x6C)
ACCENT_TEAL  = RGBColor(0x00, 0x7A, 0x87)
MED_GRAY     = RGBColor(0x99, 0x99, 0x99)
WARM_GRAY    = RGBColor(0x55, 0x55, 0x55)

DECK_PATH = r"C:\Users\75565\OneDrive - Bain\Documents\GitHub\satty001\AI_in_Biopharma_Bain_POV.pptx"
prs = Presentation(DECK_PATH)

# ── We'll insert a new slide at position 10 (after slide 10, CRO Deep Dive) ──
# python-pptx doesn't have a native "insert at index" so we add at end, then reorder
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ═══════════════════════════════════════════════════════════
# DARK BACKGROUND
# ═══════════════════════════════════════════════════════════
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x14, 0x14, 0x1E)
bg.line.fill.background()

# ═══════════════════════════════════════════════════════════
# TITLE AREA
# ═══════════════════════════════════════════════════════════
# Red accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = BAIN_RED
bar.line.fill.background()

# Title
t = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(10), Inches(0.5))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "The Clinical Development Machine of the Future"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Calibri"

# Subtitle
t2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(10), Inches(0.35))
tf2 = t2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "From an 18-month sequential process to an AI-native parallel execution engine"
p2.font.size = Pt(12)
p2.font.italic = True
p2.font.color.rgb = BAIN_LIGHT
p2.font.name = "Calibri"

# ═══════════════════════════════════════════════════════════
# STAGE DATA
# ═══════════════════════════════════════════════════════════
stages = [
    {
        "name": "Protocol\nDesign",
        "icon": "01",
        "today": "Manual literature review\n6\u201312 mo development\n30%+ amendment rate",
        "future": "AI-generated protocols\nPredictive amendment\navoidance\nSynthetic control arms",
        "metric": "50%",
        "metric_label": "faster protocol\nfinalization",
    },
    {
        "name": "Site Selection\n& Activation",
        "icon": "02",
        "today": "Historical bias\nManual feasibility\n4\u20136 mo activation",
        "future": "ML scoring from\n200+ variables\nPredictive enrollment\nmodeling",
        "metric": "40%",
        "metric_label": "faster site\nactivation",
    },
    {
        "name": "Patient\nRecruitment",
        "icon": "03",
        "today": "20\u201340% screen failure\nSlow enrollment\nLimited diversity",
        "future": "AI patient matching\nfrom EHR/claims\nDigital biomarker\npre-screening",
        "metric": "2\u00d7",
        "metric_label": "enrollment\nvelocity",
    },
    {
        "name": "Trial Conduct\n& Monitoring",
        "icon": "04",
        "today": "Manual monitoring\nReactive risk mgmt\nPaper-heavy process",
        "future": "AI CRA co-pilot\nRisk-based monitoring\nReal-time deviation\ndetection",
        "metric": "60%",
        "metric_label": "fewer monitoring\nvisits needed",
    },
    {
        "name": "Data &\nSafety",
        "icon": "05",
        "today": "Manual data cleaning\nDelayed safety signals\nSiloed pharmacovigilance",
        "future": "Automated validation\nLLM signal extraction\nContinuous PV from\nsocial & EHR",
        "metric": "200K+",
        "metric_label": "hours saved\nper year",
    },
    {
        "name": "Regulatory\nSubmission",
        "icon": "06",
        "today": "Manual CSR writing\nReactive query responses\nSequential submissions",
        "future": "AI-generated CSRs\nRegulatory query\nprediction\nAutomated dossiers",
        "metric": "40%",
        "metric_label": "faster\nsubmission prep",
    },
    {
        "name": "Evidence\n& Launch",
        "icon": "07",
        "today": "Episodic advisory boards\nDelayed evidence\nDisconnected from\ncommercial",
        "future": "Agentic advisory boards\nCDS-ready evidence\nReal-time RWE\ngeneration",
        "metric": "6\u201312 mo",
        "metric_label": "earlier commercial\nreadiness",
    },
]

N = len(stages)
# Layout constants
LEFT_MARGIN = Inches(0.35)
RIGHT_MARGIN = Inches(0.35)
USABLE_W = prs.slide_width - LEFT_MARGIN - RIGHT_MARGIN
COL_GAP = Inches(0.12)
COL_W = int((USABLE_W - COL_GAP * (N - 1)) / N)

HEADER_Y = Inches(1.15)
HEADER_H = Inches(0.6)
TODAY_LABEL_Y = Inches(1.85)
TODAY_Y = Inches(2.05)
TODAY_H = Inches(1.05)
ARROW_Y = Inches(3.15)
ARROW_H = Inches(0.3)
FUTURE_LABEL_Y = Inches(3.5)
FUTURE_Y = Inches(3.7)
FUTURE_H = Inches(1.2)
METRIC_Y = Inches(5.05)
METRIC_H = Inches(0.9)

# ═══════════════════════════════════════════════════════════
# "TODAY" and "AI-NATIVE" row labels on the left
# ═══════════════════════════════════════════════════════════
# Today label (rotated effect via narrow vertical textbox)
def add_row_label(slide, text, y, h, color):
    """Vertical row label on the far left."""
    t = slide.shapes.add_textbox(Inches(0.0), y, Inches(0.35), h)
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

# ═══════════════════════════════════════════════════════════
# DRAW STAGE COLUMNS
# ═══════════════════════════════════════════════════════════
for i, stage in enumerate(stages):
    x = LEFT_MARGIN + i * (COL_W + COL_GAP)

    # ── Stage header (number + name) ──
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, HEADER_Y, COL_W, HEADER_H)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = BAIN_RED
    hdr.line.fill.background()
    hdr_tf = hdr.text_frame
    hdr_tf.word_wrap = True
    hdr_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Number
    p_num = hdr_tf.paragraphs[0]
    p_num.text = stage["icon"]
    p_num.font.size = Pt(8)
    p_num.font.bold = True
    p_num.font.color.rgb = RGBColor(0xFF, 0xBD, 0xBD)
    p_num.font.name = "Calibri"
    p_num.space_after = Pt(0)
    # Name
    p_name = hdr_tf.add_paragraph()
    p_name.text = stage["name"]
    p_name.font.size = Pt(9)
    p_name.font.bold = True
    p_name.font.color.rgb = WHITE
    p_name.font.name = "Calibri"
    p_name.alignment = PP_ALIGN.CENTER
    p_name.space_before = Pt(0)

    # ── "Today" card ──
    today_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, TODAY_Y, COL_W, TODAY_H)
    today_card.fill.solid()
    today_card.fill.fore_color.rgb = DARK_CARD
    today_card.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    today_card.line.width = Pt(0.5)
    tc_tf = today_card.text_frame
    tc_tf.word_wrap = True
    tc_tf.margin_left = Inches(0.06)
    tc_tf.margin_right = Inches(0.06)
    tc_tf.margin_top = Inches(0.04)
    for j, line in enumerate(stage["today"].split("\n")):
        if j == 0:
            p = tc_tf.paragraphs[0]
        else:
            p = tc_tf.add_paragraph()
        p.text = "\u2022 " + line
        p.font.size = Pt(7.5)
        p.font.color.rgb = MED_GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(1)

    # ── Down arrow ──
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, x + COL_W // 2 - Inches(0.15), ARROW_Y,
        Inches(0.3), ARROW_H
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BAIN_LIGHT
    arrow.line.fill.background()

    # ── "AI Future" card ──
    future_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, FUTURE_Y, COL_W, FUTURE_H)
    future_card.fill.solid()
    future_card.fill.fore_color.rgb = RGBColor(0x2A, 0x0A, 0x0A)
    future_card.line.color.rgb = BAIN_RED
    future_card.line.width = Pt(1)
    fc_tf = future_card.text_frame
    fc_tf.word_wrap = True
    fc_tf.margin_left = Inches(0.06)
    fc_tf.margin_right = Inches(0.06)
    fc_tf.margin_top = Inches(0.04)
    for j, line in enumerate(stage["future"].split("\n")):
        if j == 0:
            p = fc_tf.paragraphs[0]
        else:
            p = fc_tf.add_paragraph()
        p.text = "\u2713 " + line
        p.font.size = Pt(7.5)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_after = Pt(1)

    # ── Metric badge ──
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, METRIC_Y, COL_W, METRIC_H)
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x36)
    badge.line.color.rgb = ACCENT_GREEN
    badge.line.width = Pt(0.75)
    b_tf = badge.text_frame
    b_tf.word_wrap = True
    b_tf.margin_left = Inches(0.04)
    b_tf.margin_right = Inches(0.04)
    b_tf.margin_top = Inches(0.04)
    # Big number
    p_m = b_tf.paragraphs[0]
    p_m.text = stage["metric"]
    p_m.font.size = Pt(18)
    p_m.font.bold = True
    p_m.font.color.rgb = ACCENT_GREEN
    p_m.font.name = "Calibri"
    p_m.alignment = PP_ALIGN.CENTER
    p_m.space_after = Pt(0)
    # Label
    p_ml = b_tf.add_paragraph()
    p_ml.text = stage["metric_label"]
    p_ml.font.size = Pt(7)
    p_ml.font.color.rgb = RGBColor(0x80, 0xCC, 0xB4)
    p_ml.font.name = "Calibri"
    p_ml.alignment = PP_ALIGN.CENTER
    p_ml.space_before = Pt(0)

    # ── Connecting arrow to next stage (horizontal) ──
    if i < N - 1:
        arr_x = x + COL_W + Inches(0.01)
        arr_y = HEADER_Y + HEADER_H // 2 - Inches(0.06)
        conn = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arr_x, arr_y, COL_GAP - Inches(0.02), Inches(0.12)
        )
        conn.fill.solid()
        conn.fill.fore_color.rgb = RGBColor(0xFF, 0x88, 0x88)
        conn.line.fill.background()

# ═══════════════════════════════════════════════════════════
# ROW LABELS
# ═══════════════════════════════════════════════════════════
# "TODAY" label
tl = slide.shapes.add_textbox(Inches(0.35), TODAY_Y - Inches(0.18), Inches(1.5), Inches(0.18))
tf_tl = tl.text_frame
p_tl = tf_tl.paragraphs[0]
p_tl.text = "TODAY"
p_tl.font.size = Pt(7)
p_tl.font.bold = True
p_tl.font.color.rgb = MED_GRAY
p_tl.font.name = "Calibri"

# "AI-NATIVE FUTURE" label
fl = slide.shapes.add_textbox(Inches(0.35), FUTURE_Y - Inches(0.18), Inches(1.5), Inches(0.18))
tf_fl = fl.text_frame
p_fl = tf_fl.paragraphs[0]
p_fl.text = "AI-NATIVE FUTURE"
p_fl.font.size = Pt(7)
p_fl.font.bold = True
p_fl.font.color.rgb = BAIN_LIGHT
p_fl.font.name = "Calibri"

# "IMPACT" label
il = slide.shapes.add_textbox(Inches(0.35), METRIC_Y - Inches(0.18), Inches(1.5), Inches(0.18))
tf_il = il.text_frame
p_il = tf_il.paragraphs[0]
p_il.text = "IMPACT"
p_il.font.size = Pt(7)
p_il.font.bold = True
p_il.font.color.rgb = ACCENT_GREEN
p_il.font.name = "Calibri"

# ═══════════════════════════════════════════════════════════
# BOTTOM SUMMARY BANNER
# ═══════════════════════════════════════════════════════════
banner_y = Inches(6.15)
banner = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), banner_y,
    prs.slide_width - Inches(0.7), Inches(0.55)
)
banner.fill.solid()
banner.fill.fore_color.rgb = BAIN_DARK
banner.line.fill.background()
b_tf = banner.text_frame
b_tf.word_wrap = True
p_b = b_tf.paragraphs[0]
p_b.text = "NET IMPACT:   30\u201340% reduction in trial timelines   |   25% lower cost per patient   |   Earlier, stronger evidence packages   |   Continuous safety signal detection"
p_b.font.size = Pt(11)
p_b.font.bold = True
p_b.font.color.rgb = WHITE
p_b.font.name = "Calibri"
p_b.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════
# HORIZONTAL TIMELINE ARROW (behind the header cards)
# ═══════════════════════════════════════════════════════════
# A thin line under the headers to show flow
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, LEFT_MARGIN, HEADER_Y + HEADER_H + Inches(0.01),
    USABLE_W, Inches(0.02)
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x55, 0x33, 0x33)
line.line.fill.background()

# ═══════════════════════════════════════════════════════════
# FOOTER
# ═══════════════════════════════════════════════════════════
ft = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.4), Inches(8), Inches(0.3))
tf_ft = ft.text_frame
p_ft = tf_ft.paragraphs[0]
p_ft.text = "Bain & Company  |  Confidential  |  February 2026"
p_ft.font.size = Pt(8)
p_ft.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p_ft.font.name = "Calibri"

# Page number
pn = slide.shapes.add_textbox(prs.slide_width - Inches(1), prs.slide_height - Inches(0.4), Inches(0.5), Inches(0.3))
tf_pn = pn.text_frame
p_pn = tf_pn.paragraphs[0]
p_pn.text = "11"
p_pn.font.size = Pt(8)
p_pn.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p_pn.font.name = "Calibri"
p_pn.alignment = PP_ALIGN.RIGHT

# ═══════════════════════════════════════════════════════════
# REORDER: Move the new slide to position 10 (index 10, after CRO Deep Dive)
# ═══════════════════════════════════════════════════════════
# The slide was appended at the end. We need to move it to index 10.
from lxml import etree

prs_element = prs.part._element
ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
sldIdLst = prs_element.find(f'{ns}sldIdLst')
sldId_elements = list(sldIdLst)
# The new slide is the last one
new_sldId = sldId_elements[-1]
# Remove from current position
sldIdLst.remove(new_sldId)
# Insert at position 10 (after slide 10 which is index 9)
target_index = 10  # 0-based: after slide 10 (CRO Deep Dive is index 9)
children = list(sldIdLst)
if target_index >= len(children):
    sldIdLst.append(new_sldId)
else:
    children[target_index].addprevious(new_sldId)

prs.save(DECK_PATH)
print(f"Diagram slide inserted at position 11 (after CRO Deep Dive)")
print(f"Total slides: {len(prs.slides)}")
