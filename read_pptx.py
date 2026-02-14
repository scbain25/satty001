"""Extract all text from a .pptx file."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches

path = sys.argv[1]
prs = Presentation(path)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*80}")
    print(f"SLIDE {i}")
    print(f"{'='*80}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(text)
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(" | ".join(cells))
    # Check for notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame
        notes_text = notes.text.strip()
        if notes_text:
            print(f"\n[SPEAKER NOTES]: {notes_text}")
